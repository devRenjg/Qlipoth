from pathlib import Path
import re
from docx import Document as DocxDocument
from python_calamine import CalamineWorkbook
from pptx import Presentation


def _looks_like_noise_id(name: str) -> bool:
    """判断 @后内容是否为噪声(企微协作残片/内部ID)，而非真实人名。"""
    # 纯数字、纯符号、含数字的短串
    if re.fullmatch(r"[\d_]+", name):
        return True
    if re.fullmatch(r"[A-Za-z0-9_]+", name):
        if len(name) <= 3:
            return True
        # 企微协作残片：auto/jauto 前缀、wingding、纯大写+数字混合(非人名)
        if re.search(r"auto[Z0-9]|wingding|jauto", name, re.I):
            return True
        if re.fullmatch(r"[A-Z0-9]{4,}", name):              # ZHCBCDD1Zj / FF0000 全大写数字串
            return True
        # 字母+J 开头的残片：eJ/oJ/nJ/fJ/sJ/tJ/vJ/jJ + 任意(含小写 oJll/fJop)
        if re.match(r"[a-z]J[A-Za-z0-9]", name) or re.match(r"(e[JT]|[onfsvtj]J)", name):
            return True
        # j[大写/数字混合]j 残片：jG0Bj / jI1Dj / jW8Rj 等
        if re.fullmatch(r"j[A-Z0-9][0-9A-Z][A-Za-z0-9]*j?", name):
            return True
        # 含大写连串+数字的随机串(ZHCBCDD1Zj 类)，排除正常驼峰人名(首字母大写其余小写)
        if re.search(r"[A-Z]{3,}", name) and re.search(r"[0-9]", name):
            return True
        if re.search(r"P[A-Z0-9]{2,}", name):                # xxPECD / PB25 残片
            return True
        if re.fullmatch(r"[a-z0-9]{1,4}", name):             # 纯小写短串
            return True
        # 纯小写字母+数字混合的随机串(rzm5v3zp / gf4r6s / rz5ecu7r4u 类残片)
        if re.fullmatch(r"[a-z0-9]+", name) and re.search(r"[0-9]", name):
            return True
        # 纯小写字母随机串残片:rz开头(企微残片特征) 或 元音比过低(辅音堆叠，非真实人名)
        if re.fullmatch(r"[a-z]{5,}", name):
            if name.startswith("rz"):
                return True
            vowel = sum(1 for c in name if c in "aeiou") / len(name)
            if vowel < 0.25:
                return True
    return False


# @后跟的明显是正文短语而非人名的词(误抓)
_NON_NAME_HINT = re.compile(r"(无法识别|直播间|研发负责人|确认|后台|观看人数|精彩|资源|页面|开启|系统|平台$|项目$)")


def extract_owners(text: str) -> list[str]:
    """从文档中提取 @提及 的真实人名作为负责人，过滤企微协作残片/内部ID/正文误抓。"""
    matches = re.findall(r"@([一-鿿\w]{2,10})", text)
    seen = set()
    owners = []
    for name in matches:
        if name in seen:
            continue
        seen.add(name)
        if _looks_like_noise_id(name):          # 噪声ID残片
            continue
        if _NON_NAME_HINT.search(name):         # @后跟的是正文短语
            continue
        owners.append(name)
    return owners


def parse_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    lines = []
    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1
            lines.append(f"{'#' * level} {para.text}")
        elif para.text.strip():
            lines.append(para.text)
    for table in doc.tables:
        lines.append("")
        headers = [cell.text.strip() for cell in table.rows[0].cells]
        lines.append("| " + " | ".join(headers) + " |")
        lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for row in table.rows[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
        lines.append("")
    return "\n".join(lines)


def parse_xlsx(file_path: str) -> str:
    wb = CalamineWorkbook.from_path(file_path)
    lines = []
    for sheet_name in wb.sheet_names:
        data = wb.get_sheet_by_name(sheet_name).to_python()
        if not data:
            continue
        lines.append(f"## {sheet_name}")
        lines.append("")

        max_col = max(len(row) for row in data)
        normalized = []
        for row in data:
            cells = []
            for val in row:
                if val is None:
                    cells.append("")
                else:
                    s = str(val).replace("\n", " ").replace("\r", " ").replace("|", "\\|")
                    cells.append(s)
            while len(cells) < max_col:
                cells.append("")
            normalized.append(cells)

        groups = _split_column_groups(normalized)
        if groups:
            for group in groups:
                lines.extend(_render_table(group))
                lines.append("")
        else:
            lines.extend(_render_table(normalized))
            lines.append("")
    return "\n".join(lines)


def _split_column_groups(rows: list[list[str]]) -> list[list[list[str]]] | None:
    """Split a wide table into groups separated by empty columns."""
    if not rows or len(rows[0]) <= 6:
        return None

    col_count = len(rows[0])
    empty_cols = []
    for col_idx in range(col_count):
        if all(row[col_idx].strip() == "" for row in rows):
            empty_cols.append(col_idx)

    if not empty_cols:
        return None

    boundaries = []
    start = 0
    for ec in empty_cols:
        if ec > start:
            boundaries.append((start, ec))
        start = ec + 1
    if start < col_count:
        boundaries.append((start, col_count))

    if len(boundaries) <= 1:
        return None

    groups = []
    for col_start, col_end in boundaries:
        group = []
        for row in rows:
            group_row = row[col_start:col_end]
            if any(cell.strip() for cell in group_row):
                group.append(group_row)
        if group:
            groups.append(group)
    return groups


def _render_table(rows: list[list[str]]) -> list[str]:
    """Render rows as a markdown table."""
    if not rows:
        return []
    col_count = max(len(row) for row in rows)
    lines = []
    headers = rows[0]
    while len(headers) < col_count:
        headers.append("")
    lines.append("| " + " | ".join(headers) + " |")
    lines.append("| " + " | ".join(["---"] * col_count) + " |")
    for row in rows[1:]:
        while len(row) < col_count:
            row.append("")
        lines.append("| " + " | ".join(row) + " |")
    return lines


def parse_xls(file_path: str) -> str:
    return parse_xlsx(file_path)


def parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
        lines.append("")
    return "\n".join(lines)


def parse_text(file_path: str) -> str:
    return Path(file_path).read_text(encoding="utf-8")


def parse_pdf(file_path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    lines = []
    has_text = False
    for i, page in enumerate(reader.pages, 1):
        try:
            text = page.extract_text()
        except Exception:
            text = ""
        if text and text.strip():
            lines.append(f"## 第 {i} 页")
            lines.append(text.strip())
            lines.append("")
            has_text = True

    if has_text:
        return "\n".join(lines)

    # Image-based PDF: render each page to an image via PyMuPDF, then OCR.
    # PyMuPDF rendering avoids pypdf's decompression limit on embedded images.
    return _ocr_pdf_pages(file_path)


def _ocr_image_tesseract(img, tesseract_path: str) -> str:
    """OCR a PIL image, slicing tall images into strips.

    Tesseract rejects images above ~100M pixels ("Image too large"), which
    happens with long marketing/poster PDFs. Slice by height and concatenate.
    """
    import subprocess
    import tempfile

    W, H = img.size
    strip_height = 2000
    parts = []
    for top in range(0, H, strip_height):
        crop = img.crop((0, top, W, min(top + strip_height, H)))
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                tmp_path = tmp.name
            crop.save(tmp_path)
            proc = subprocess.run(
                [tesseract_path, tmp_path, "stdout", "-l", "chi_sim+eng", "--psm", "6"],
                capture_output=True, text=True, timeout=60, encoding="utf-8"
            )
            text = (proc.stdout or "").strip()
            if text:
                parts.append(text)
        except Exception:
            continue
        finally:
            if tmp_path:
                Path(tmp_path).unlink(missing_ok=True)
    return "\n".join(parts)


def _ocr_pdf_pages(file_path: str, max_pages: int = 30) -> str:
    """Render each PDF page to an image with PyMuPDF, then OCR via Tesseract.

    Rendering whole pages avoids pypdf's per-image decompression limit, which
    fails on PDFs containing large embedded images (LimitReachedError).
    """
    import fitz
    from PIL import Image

    Image.MAX_IMAGE_PIXELS = None  # allow tall poster-style pages
    tesseract_path = "C:/Program Files/Tesseract-OCR/tesseract.exe"

    results = []
    try:
        doc = fitz.open(file_path)
    except Exception:
        return ""

    try:
        for i in range(min(doc.page_count, max_pages)):
            try:
                page = doc.load_page(i)
                # 2x zoom for sharper text -> better OCR accuracy
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text = _ocr_image_tesseract(img, tesseract_path)
                if text:
                    results.append(f"## 第 {i+1} 页\n\n{text}")
            except Exception:
                continue
    finally:
        doc.close()

    return "\n\n".join(results)


PARSERS = {
    ".docx": parse_docx,
    ".xlsx": parse_xlsx,
    ".xls": parse_xls,
    ".pptx": parse_pptx,
    ".pdf": parse_pdf,
    ".txt": parse_text,
    ".md": parse_text,
}


def parse_file(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    parser = PARSERS.get(ext)
    if not parser:
        raise ValueError(f"不支持的文件格式: {ext}")
    return parser(file_path)
